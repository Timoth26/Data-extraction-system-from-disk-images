from transformers import pipeline
import pdfplumber
from docx import Document
from bs4 import BeautifulSoup
from lxml import etree
from collections import Counter
from ocr import ocr
import csv
import json
import pptx
import zipfile
import sqlite3
from email import message_from_file
import ebooklib

import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from collections import Counter
from transformers import pipeline


def analyze_file(file_path, ner_pipeline, score_threshold):
    try:
        print(f"[INFO] Starting analysis of file: {file_path}")
        if file_path.lower().endswith(".pdf"):
            print(f"[INFO] Extracting text from PDF: {file_path}")
            text = extract_text_from_pdf(file_path)
        elif file_path.lower().endswith((".txt", ".log", ".eml")):
            print(f"[INFO] Reading text from TXT/LOG/EML: {file_path}")
            with open(file_path, "r", encoding="utf-8") as file:
                text = file.read()
        elif file_path.lower().endswith(".docx"):
            print(f"[INFO] Extracting text from DOCX: {file_path}")
            text = extract_text_from_docx(file_path)
        elif file_path.lower().endswith((".html", ".xml")):
            print(f"[INFO] Extracting text from HTML/XML: {file_path}")
            text = extract_text_from_html(file_path)
        elif file_path.lower().endswith(".csv"):
            print(f"[INFO] Extracting text from CSV: {file_path}")
            text = extract_text_from_csv(file_path)
        elif file_path.lower().endswith(".json"):
            print(f"[INFO] Extracting text from JSON: {file_path}")
            text = extract_text_from_json(file_path)
        elif file_path.lower().endswith(".pptx"):
            print(f"[INFO] Extracting text from PPTX: {file_path}")
            text = extract_text_from_pptx(file_path)
        elif file_path.lower().endswith(".odt"):
            print(f"[INFO] Extracting text from ODT: {file_path}")
            text = extract_text_from_odt(file_path)
        elif file_path.lower().endswith(".md"):
            print(f"[INFO] Extracting text from Markdown: {file_path}")
            text = extract_text_from_md(file_path)
        elif file_path.lower().endswith(".msg"):
            print(f"[INFO] Extracting text from MSG: {file_path}")
            text = extract_text_from_msg(file_path)
        elif file_path.lower().endswith(".epub"):
            print(f"[INFO] Extracting text from EPUB: {file_path}")
            text = extract_text_from_epub(file_path)
        elif file_path.lower().endswith(".db"):
            print(f"[INFO] Extracting text from Database: {file_path}")
            text = extract_text_from_db(file_path)
        elif file_path.lower().endswith((".png", ".jpeg", ".jpg")):
            print(f"[INFO] Performing OCR on Image: {file_path}")
            text = ocr(file_path)
        else:
            print(f"[WARNING] Unsupported file format: {file_path}")
            return file_path, None  # Unsupported file format

        print(f"[INFO] Running NER model on file: {file_path}")
        ner_results = ner_pipeline(text)
        filtered_results = [entity for entity in ner_results if entity['score'] >= score_threshold]
        print(f"[INFO] Completed NER for file: {file_path} ({len(filtered_results)} entities detected)")
        return file_path, filtered_results

    except Exception as e:
        print(f"[ERROR] Exception during analysis of file {file_path}: {e}")

def analyze_files(file_paths, output_file, score_threshold=0.90, max_workers=4):
    ner_pipeline = pipeline("token-classification", model='lakshyakh93/deberta_finetuned_pii', device=-1, aggregation_strategy="simple")
    results = {}

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {executor.submit(analyze_file, file_path, ner_pipeline, score_threshold): file_path for file_path in file_paths}
        with open(output_file, "w", encoding="utf-8") as result_file:
            for future in as_completed(future_to_file):
                file_path = future_to_file[future]
                try:
                    file_path, filtered_results = future.result()
                    results[file_path] = filtered_results

                    result_file.write(f"Analysis of file: {file_path}\n\n")
                    if isinstance(filtered_results, list) and filtered_results:
                        for entity in filtered_results:
                            result_file.write(f"Entity: {entity['word']}\nType: {entity['entity_group']}\n"
                                              f"Score: {entity['score']:.4f}\nStart: {entity['start']}\nEnd: {entity['end']}\n\n")
                    elif isinstance(filtered_results, str):
                        result_file.write(f"Error: {filtered_results}\n\n")
                    else:
                        result_file.write("No entities detected above the threshold.\n\n")
                    result_file.write("-" * 50 + "\n\n")

                    print(f"[INFO] Successfully wrote results for file: {file_path}")
                except Exception as e:
                    print(f"[ERROR] Exception while writing results for {file_path}: {e}")

    print("[INFO] Completed analysis for all files.")
    return results


def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_html(html_path):
    with open(html_path, "r", encoding="utf-8") as file:
        soup = BeautifulSoup(file, "html.parser")
        return soup.get_text()
    
def extract_text_from_xml(xml_path):
    with open(xml_path, "rb") as file:
        tree = etree.parse(file)
    return " ".join(tree.xpath("//text()"))

def extract_text_from_csv(csv_path):
    text = ""
    with open(csv_path, "r", encoding="utf-8") as file:
        reader = csv.reader(file)
        for row in reader:
            text += " ".join(row) + "\n"
    return text

def extract_text_from_json(json_path):
    with open(json_path, "r", encoding="utf-8") as file:
        data = json.load(file)
    return json.dumps(data)

def extract_text_from_pptx(pptx_path):
    text = ""
    presentation = pptx.Presentation(pptx_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_odt(odt_path):
    text = ""
    with zipfile.ZipFile(odt_path) as zf:
        xml_content = zf.read("content.xml")
        tree = etree.XML(xml_content)
        text = " ".join(tree.xpath("//text:p//text()"))
    return text

def extract_text_from_md(md_path):
    with open(md_path, "r", encoding="utf-8") as file:
        return file.read()

def extract_text_from_msg(msg_path):
    with open(msg_path, "r", encoding="utf-8") as file:
        content = file.read()
    return content

def extract_text_from_eml(eml_path):
    with open(eml_path, "r", encoding="utf-8") as file:
        msg = message_from_file(file)
        return msg.get_payload(decode=True).decode("utf-8", errors="ignore")

def extract_text_from_epub(epub_path):
    book = ebooklib.epub.read_epub(epub_path)
    text = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            text += item.get_body_content().decode("utf-8")
    return text

def extract_text_from_db(db_path):
    text = ""
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = cursor.fetchall()
    for table in tables:
        table_name = table[0]
        cursor.execute(f"SELECT * FROM {table_name}")
        rows = cursor.fetchall()
        for row in rows:
            text += " ".join(str(col) for col in row) + "\n"
    conn.close()
    return text

def count_entities(results):
    if results:
        entity_groups = []
        for key, value in results.items():
            entity_groups.extend([item['entity_group'] for item in value])
        
        entity_group_counts = Counter(entity_groups)
        return dict(entity_group_counts)
    else:
        return None
